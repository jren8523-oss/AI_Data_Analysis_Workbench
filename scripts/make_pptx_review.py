# -*- coding: utf-8 -*-
"""《工作台制作复盘》PPT 生成器：13 页，讲"怎么用 AI 做出数据分析工作台"。

独立于 engine/ppt_out.py（那份是"交付总览"，讲工作台能干嘛）；
这份是"制作复盘"，讲工作台是怎么一步步搭出来的，给学生看"怎么用 AI 做工具"。
配色复用项目人邮蓝体系，与 ppt_out.py 一致。

用法：
    python scripts/make_pptx_review.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------- 主题色（人邮蓝体系，与 ppt_out.py 一致） ----------
BRAND = RGBColor(0x00, 0x5D, 0xA2)
BRAND_DARK = RGBColor(0x00, 0x44, 0x7A)
BRAND_LIGHT = RGBColor(0xE8, 0xF1, 0xF8)
ACCENT = RGBColor(0xF2, 0xA9, 0x00)
ACCENT_BG = RGBColor(0xFA, 0xEE, 0xDA)
ACCENT_TXT = RGBColor(0x63, 0x38, 0x06)
OK = RGBColor(0x1E, 0x8E, 0x5A)
TEXT = RGBColor(0x2D, 0x37, 0x48)
GRAY = RGBColor(0x8A, 0x94, 0xA0)
LIGHT_GRAY = RGBColor(0xF1, 0xEF, 0xE8)
DARK_GRAY = RGBColor(0x44, 0x44, 0x41)
SUBTITLE = RGBColor(0xD9, 0xE7, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)

# ---------- 截图（xhs / 推荐帖，嵌入对应页） ----------
IMG_DIR = Path(r"C:\Users\lenovo\.workbuddy\clipboard-images")
IMG_GRILL = IMG_DIR / "clipboard-2026-09-02T11-40-30-209Z-1cfe0953.png"   # grill-me 说明
IMG_REVIEW = IMG_DIR / "clipboard-2026-09-02T11-40-30-219Z-0cfae8cb.jpg"  # 审美 skill 测评总览
IMG_TASTE = IMG_DIR / "clipboard-2026-09-02T11-40-30-223Z-79d7fec1.png"   # taste-skill 实机
IMG_UIUX = IMG_DIR / "clipboard-2026-09-02T11-40-30-229Z-97785283.png"    # uiuxpromax 实机
IMG_FRONT = IMG_DIR / "clipboard-2026-09-02T11-40-30-233Z-b0a15f2d.png"   # frontend-design 实机


def new_pres() -> Presentation:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.2)
    return sp


def text(slide, l, t, w, h, s, size=14, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "微软雅黑"
    return tb


def bar(slide, title, subtitle=""):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.0))
    sp.fill.solid()
    sp.fill.fore_color.rgb = BRAND
    sp.line.fill.background()
    tf = sp.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "微软雅黑"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = SUBTITLE
        r2.font.name = "微软雅黑"


def pic(slide, path, l, t, h=None, w=None):
    path = Path(path)
    if not path.exists():
        return
    kw = {}
    if h:
        kw["height"] = Inches(h)
    if w:
        kw["width"] = Inches(w)
    slide.shapes.add_picture(str(path), Inches(l), Inches(t), **kw)


def build() -> Path:
    prs = new_pres()

    # ============ P1 封面 ============
    s = blank(prs)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.35))
    top.fill.solid(); top.fill.fore_color.rgb = BRAND; top.line.fill.background()
    text(s, 1.0, 2.0, 11.3, 0.5, "AI WORKBENCH", 16, True, BRAND, PP_ALIGN.CENTER)
    text(s, 1.0, 2.6, 11.3, 1.0, "从零做一个数据分析工作台", 40, True, BRAND_DARK, PP_ALIGN.CENTER)
    text(s, 1.0, 3.75, 11.3, 0.5, "一条 AI 协作流水线 · 让学生看懂「怎么用 AI 做工具」", 16, color=GRAY, align=PP_ALIGN.CENTER)
    text(s, 1.0, 6.6, 11.3, 0.4, "单文件 HTML · 双击即开 · 数据不出浏览器", 12, color=GRAY, align=PP_ALIGN.CENTER)

    # ============ P2 为什么是工作台 ============
    s = blank(prs)
    bar(s, "为什么是「工作台」，不是「单次任务」？", "可复用，是工具与一次性请求的本质区别")
    rect(s, 0.7, 1.5, 5.7, 4.0, LIGHT_GRAY)
    text(s, 1.0, 1.8, 5.1, 0.5, "单次任务", 20, True, GRAY)
    text(s, 1.0, 2.55, 5.1, 2.8, "· 每次把数据 + 需求喂给 AI\n· 现算现出，结果不沉淀\n· 下次换数据，重头再来\n· 口径可能每次不一样", 15, color=TEXT)
    rect(s, 6.9, 1.5, 5.7, 4.0, BRAND_LIGHT)
    text(s, 7.2, 1.8, 5.1, 0.5, "工作台", 20, True, BRAND_DARK)
    text(s, 7.2, 2.55, 5.1, 2.8, "· 需求 / 方法 / 口径一次固化\n· 上传数据 → 点一下 → 出报告\n· 随取随用，永久复用\n· 双引擎 + QA，结果可验证", 15, color=BRAND_DARK)
    rect(s, 0.7, 5.95, 11.9, 1.0, ACCENT_BG)
    text(s, 1.0, 6.2, 11.3, 0.6, "单次是「求 AI 帮一次忙」，工作台是「把 AI 的能力装进一个随取随用的工具」", 15, True, ACCENT_TXT)

    # ============ P3 六步流水线 ============
    s = blank(prs)
    bar(s, "全流程总览：六步流水线", "「工作台搭建师」专家贯穿全程")
    steps = [
        ("1", "让 AI 读懂需求", "喂课程 PPT"),
        ("2", "召唤「工作台搭建师」", "Expert"),
        ("3", "设计对齐", "grill-me"),
        ("4", "画界面", "frontend-design"),
        ("5", "搭 UI + 填功能", "Codex CLI"),
        ("6", "质量保障", "karpathy / superpower"),
    ]
    for i, (num, name, tool) in enumerate(steps):
        x = 0.7 + (i % 3) * 4.05
        y = 1.7 + (i // 3) * 2.4
        rect(s, x, y, 3.75, 2.0, BRAND_LIGHT if i % 2 == 0 else LIGHT_GRAY)
        text(s, x + 0.25, y + 0.25, 0.9, 0.9, num, 34, True, ACCENT)
        text(s, x + 1.05, y + 0.3, 2.55, 0.9, name, 15, True, BRAND_DARK)
        text(s, x + 1.05, y + 1.3, 2.55, 0.5, tool, 12, color=GRAY)
    text(s, 0.7, 6.85, 11.9, 0.4, "专家管「怎么想」，技能管「会不会做」——一个专家贯穿，五把技能按环节换", 12, color=GRAY, align=PP_ALIGN.CENTER)

    # ============ P4 第一步 读懂需求 ============
    s = blank(prs)
    bar(s, "第一步 · 让 AI 读懂需求", "需求不是凭空提，是从课程材料里「喂」出来的")
    flow = [("课程 PPT", "《统计与数据分析基础》"), ("AI 总结", "提炼实训与统计口径"), ("表述需求", "我要做一个数据分析工作台")]
    for i, (t, d) in enumerate(flow):
        x = 0.8 + i * 4.1
        rect(s, x, 2.2, 3.5, 1.9, BRAND_LIGHT)
        text(s, x + 0.2, 2.55, 3.1, 0.6, t, 18, True, BRAND_DARK, PP_ALIGN.CENTER)
        text(s, x + 0.2, 3.3, 3.1, 0.7, d, 12, color=TEXT, align=PP_ALIGN.CENTER)
        if i < 2:
            text(s, x + 3.5, 2.75, 0.6, 0.6, "→", 26, True, ACCENT, PP_ALIGN.CENTER)
    rect(s, 0.8, 4.9, 11.7, 1.5, LIGHT_GRAY)
    text(s, 1.1, 5.2, 11.1, 1.0, "要点：先把课程材料交给 AI 让它「读懂」，再用一句话把需求说清楚。\n先总结、再定义——避免一上来就让它瞎猜。", 14, color=TEXT)

    # ============ P5 第二步 召唤专家 ============
    s = blank(prs)
    bar(s, "第二步 · 召唤「工作台搭建师」专家", "全程以搭建者的身份和视角推进")
    rect(s, 0.8, 1.7, 5.7, 2.6, BRAND_LIGHT)
    text(s, 1.1, 2.0, 5.1, 0.6, "Expert 专家", 20, True, BRAND_DARK)
    text(s, 1.1, 2.75, 5.1, 1.4, "管「怎么想」\n定义身份、角色和思维方式\n同一时间只认一个", 14, color=TEXT)
    rect(s, 6.9, 1.7, 5.7, 2.6, LIGHT_GRAY)
    text(s, 7.2, 2.0, 5.1, 0.6, "Skill 技能", 20, True, DARK_GRAY)
    text(s, 7.2, 2.75, 5.1, 1.4, "管「会不会做」\n装进新能力，按需开关\n可以同时挂好几个", 14, color=TEXT)
    rect(s, 0.8, 4.7, 11.8, 1.7, ACCENT_BG)
    text(s, 1.1, 5.0, 11.2, 1.2, "最佳搭配 = 专家 + 技能\n开「工作台搭建师」定方向，再按环节挂 grill-me / frontend-design 等技能", 15, True, ACCENT_TXT)

    # ============ P6 第三步 设计对齐（grill-me） ============
    s = blank(prs)
    bar(s, "第三步 · 设计对齐", "用 grill-me 把需求颗粒度问清楚")
    qa = [("功能", "逐层追问，把模糊需求问成清晰计划"),
          ("为什么", "没有它，AI 只能瞎猜你的真实意图"),
          ("不用会怎样", "口径没人追问 → 双引擎对不上，改两遍")]
    for i, (t, d) in enumerate(qa):
        y = 1.7 + i * 1.5
        rect(s, 0.7, y, 6.6, 1.3, BRAND_LIGHT if i % 2 == 0 else LIGHT_GRAY)
        text(s, 0.95, y + 0.3, 1.9, 0.9, t, 16, True, BRAND_DARK)
        text(s, 2.9, y + 0.18, 4.3, 1.0, d, 12, color=TEXT)
    pic(s, IMG_GRILL, 7.7, 1.6, h=4.6)
    text(s, 7.7, 6.35, 4.9, 0.4, "▲ grill-me：访谈式提问，帮你把构思变清晰", 11, color=GRAY)

    # ============ P7 第四步 画界面·为什么用设计 skill ============
    s = blank(prs)
    bar(s, "第四步 · 画界面，为什么要用设计 skill？", "通用 AI 出「能跑但丑」，设计 skill 出「有审美」")
    rect(s, 0.7, 1.5, 4.0, 2.2, LIGHT_GRAY)
    text(s, 0.95, 1.78, 3.5, 0.5, "通用 AI 直接写", 15, True, GRAY)
    text(s, 0.95, 2.4, 3.5, 1.2, "烂大街渐变、AI 味重\n能跑，但没法直接投影", 12, color=TEXT)
    rect(s, 0.7, 3.95, 4.0, 2.2, BRAND_LIGHT)
    text(s, 0.95, 4.23, 3.5, 0.5, "设计 skill 来画", 15, True, BRAND_DARK)
    text(s, 0.95, 4.85, 3.5, 1.2, "有自己的审美、极简细节控\n生成 HTML，对话即可微调", 12, color=BRAND_DARK)
    pic(s, IMG_REVIEW, 5.05, 1.5, h=4.6)
    text(s, 5.05, 6.25, 3.2, 0.4, "▲ 5 个审美 skill 测评", 11, color=GRAY)
    rect(s, 8.6, 1.7, 4.0, 4.4, ACCENT_BG)
    text(s, 8.85, 2.1, 3.5, 3.8, "结论\n\n5 个审美 skill 实测，\n这套工作台选\nfrontend-design", 16, True, ACCENT_TXT)

    # ============ P8 第四步 画界面·实机对比 ============
    s = blank(prs)
    bar(s, "第四步 · 5 个审美 skill 实机对比", "同一套设计 prompt，不同 skill 出的页面")
    cards = [(IMG_TASTE, "taste-skill", "英文极简 portfolio，做中文教学工具水土不服"),
             (IMG_UIUX, "UI-UX-PRO-MAX", "中文能打，但实测略有 AI 味"),
             (IMG_FRONT, "frontend-design", "中文质感最好，极简细节控，胜出")]
    for i, (img, name, desc) in enumerate(cards):
        x = 0.7 + i * 4.1
        pic(s, img, x, 1.5, h=3.1)
        text(s, x, 4.75, 3.8, 0.5, name, 14, True, BRAND_DARK if i == 2 else TEXT)
        text(s, x, 5.25, 3.8, 0.9, desc, 11, color=TEXT)
    rect(s, 0.7, 6.5, 11.9, 0.8, BRAND)
    text(s, 0.95, 6.66, 11.4, 0.5, "对「中文 + 教学投影」场景，frontend-design 最合——三张实机 + 博主实测得出", 13, True, WHITE)

    # ============ P9 第五步 搭 UI + 填功能 ============
    s = blank(prs)
    bar(s, "第五步 · 搭 UI + 填功能", "先骨架、再微调、后填血肉")
    flow = [("搭骨架", "把三栏界面结构搭出来"), ("微调", "Codex CLI 里顺手调细节"), ("填血肉", "让 AI 填 12 项分析能力")]
    for i, (t, d) in enumerate(flow):
        x = 0.8 + i * 4.1
        rect(s, x, 2.2, 3.5, 1.9, BRAND_LIGHT)
        text(s, x + 0.2, 2.55, 3.1, 0.6, t, 18, True, BRAND_DARK, PP_ALIGN.CENTER)
        text(s, x + 0.2, 3.3, 3.1, 0.7, d, 12, color=TEXT, align=PP_ALIGN.CENTER)
        if i < 2:
            text(s, x + 3.5, 2.75, 0.6, 0.6, "→", 26, True, ACCENT, PP_ALIGN.CENTER)
    rect(s, 0.8, 4.9, 11.7, 1.5, LIGHT_GRAY)
    text(s, 1.1, 5.2, 11.1, 1.0, "要点：界面先有「形」，AI 才知道往哪填功能。\nCodex CLI 让「边聊边改代码」成为可能，微调不再靠手敲。", 14, color=TEXT)

    # ============ P10 第六步 质量保障 ============
    s = blank(prs)
    bar(s, "第六步 · 工程质量保障", "用技能守住质量，不靠肉眼")
    qa = [("功能", "karpathy / superpower 把工程质量程序化"),
          ("为什么", "AI 也会犯错，得有「验证」兜底"),
          ("不用会怎样", "众数取最小、指数自洽这种边界 bug 没人抓")]
    for i, (t, d) in enumerate(qa):
        y = 1.7 + i * 1.4
        rect(s, 0.7, y, 7.2, 1.2, BRAND_LIGHT if i % 2 == 0 else LIGHT_GRAY)
        text(s, 0.95, y + 0.28, 2.0, 0.9, t, 16, True, BRAND_DARK)
        text(s, 3.0, y + 0.16, 4.7, 1.0, d, 12, color=TEXT)
    rect(s, 8.3, 1.7, 4.3, 3.9, ACCENT_BG)
    text(s, 8.55, 2.3, 3.8, 3.0, "冒烟测试\n61 项\n\nQA 三层校验\n78 项全绿", 20, True, ACCENT_TXT, PP_ALIGN.CENTER)

    # ============ P11 成果 12 能力 ============
    s = blank(prs)
    bar(s, "成果 · 12 项分析能力", "10 个课堂实训全跑通，78 项 QA 全绿")
    caps = ["描述统计", "多指标对比", "数据清洗", "抽样估计", "统计指数", "相关回归",
            "时间序列", "分组汇总", "t 检验", "方差分析", "卡方检验", "非参数检验"]
    for i, c in enumerate(caps):
        x = 0.7 + (i % 4) * 3.05
        y = 1.7 + (i // 4) * 1.5
        rect(s, x, y, 2.8, 1.2, BRAND_LIGHT if i % 2 == 0 else LIGHT_GRAY)
        text(s, x, y + 0.38, 2.8, 0.5, c, 14, True, BRAND_DARK, PP_ALIGN.CENTER)
    rect(s, 0.7, 6.4, 11.9, 0.8, OK)
    text(s, 0.95, 6.56, 11.4, 0.5, "描述 + 推断统计全覆盖 · 双引擎口径一致 · 数据不满足条件会明确提示而非伪造", 13, True, WHITE)

    # ============ P12 给学生的一句话 ============
    s = blank(prs)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.35))
    top.fill.solid(); top.fill.fore_color.rgb = BRAND; top.line.fill.background()
    text(s, 1.0, 2.3, 11.3, 1.2, "做工具 = 选对专家 + 挂对技能 + 阶段验收", 30, True, BRAND_DARK, PP_ALIGN.CENTER)
    text(s, 1.0, 3.65, 11.3, 0.6, "而不是「一键生成」", 20, color=GRAY, align=PP_ALIGN.CENTER)
    rect(s, 2.5, 4.75, 8.3, 1.4, ACCENT_BG)
    text(s, 2.8, 5.05, 7.7, 0.9, "清晰的需求 + 趁手的技能 + 每一步都验证\n——这才是用 AI 做出可用工具的正确姿势", 15, True, ACCENT_TXT, PP_ALIGN.CENTER)

    # ============ P13 附录 术语表 ============
    s = blank(prs)
    bar(s, "附录 · 术语表", "正文用到的词，这里都有解释")
    terms = [
        ("Expert 专家", "管「怎么想」：定身份、角色和视角，一次只认一个"),
        ("Skill 技能", "管「会不会做」：装进新能力，可同时挂多个、按需开关"),
        ("grill-me", "访谈式追问，把模糊需求问成清晰可执行的计划"),
        ("frontend-design", "有审美的网页生成技能，出 HTML、响应式、可对话微调"),
        ("Codex CLI", "命令行 AI 编码工具，边聊边改代码、顺手微调"),
        ("karpathy / superpower", "把工程质量（代码质量、测试验收）程序化的技能"),
        ("双引擎", "浏览器 JS 现场算 + Python 当基准验算，口径必须一致"),
    ]
    for i, (t, d) in enumerate(terms):
        y = 1.55 + i * 0.78
        rect(s, 0.7, y, 11.9, 0.66, BRAND_LIGHT if i % 2 == 0 else LIGHT_GRAY)
        text(s, 0.95, y + 0.12, 3.4, 0.5, t, 13, True, BRAND_DARK)
        text(s, 4.5, y + 0.14, 7.9, 0.5, d, 12, color=TEXT)

    out = Path("output/工作台制作复盘.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out, len(prs.slides._sldIdLst)


if __name__ == "__main__":
    path, n = build()
    print(f"已生成: {path}（共 {n} 页）")

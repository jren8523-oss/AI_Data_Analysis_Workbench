# -*- coding: utf-8 -*-
"""交付 PPT 导出器：把工作台分析结果导出为可编辑 PPTX（独立于课程 PPT）。

分析结果（图表、标题、QA 结果）可一键打包成 PPTX，老师可直接在 PowerPoint/WPS 中
编辑与二次排版。数据全部来自 output/（build_log.json + 各任务 charts/*.png），
不重算任何指标——保证 PPT 与报告页数字一致。

用法：
    from engine import ppt_out
    ppt_out.summary_pptx("output/交付总览.pptx")          # 总览：每任务一页
    ppt_out.export_task(7, "output/task07_交付页.pptx")   # 单任务页
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from . import config

# ---------- 主题色（与工作台一致：人邮蓝） ----------
BRAND = RGBColor(0x00, 0x5D, 0xA2)
BRAND_DARK = RGBColor(0x00, 0x44, 0x7A)
BRAND_LIGHT = RGBColor(0xE8, 0xF1, 0xF8)
ACCENT = RGBColor(0xF2, 0xA9, 0x00)
OK = RGBColor(0x1E, 0x8E, 0x5A)
BAD = RGBColor(0xC0, 0x39, 0x2B)
TEXT = RGBColor(0x2D, 0x37, 0x48)
GRAY = RGBColor(0x8A, 0x94, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _new_pres() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs) -> object:
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _box(slide, l, t, w, h, text, size=14, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "微软雅黑"
    return tb


def _bar(slide, text: str, subtitle: str = ""):
    """顶部蓝色标题条（对齐参考 PPT 的色块标题）。"""
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.0))  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "微软雅黑"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(0xD9, 0xE7, 0xF4)
        r2.font.name = "微软雅黑"


def _qa_line(status: str, passed: int, checks: int) -> str:
    if status == "PASS":
        return f"✅ QA 验证：{passed}/{checks} 项检查全部通过"
    if status == "ERROR":
        return f"💥 构建出错（详见 build_log.json）"
    if status == "SKIP":
        return "⛔ 占位任务：需外部采集工具，不参与构建"
    return f"❌ QA 验证：{passed}/{checks} 通过（有失败项）"


def _task_slide(prs, task_no: int, entry: dict):
    """每任务一页：标题条 + 左侧信息 + 右侧首张图表。"""
    meta = config.TASKS[task_no]
    slide = _blank(prs)
    _bar(slide, f"任务 {task_no:02d} · {meta['title']}", meta["subtitle"])

    status = entry.get("status", "?")
    passed = entry.get("passed", 0)
    checks = entry.get("checks", 0)

    # 左侧信息
    _box(slide, 0.55, 1.45, 5.9, 0.5, meta["chapter"], 15, bold=True, color=BRAND)
    _box(slide, 0.55, 2.05, 5.9, 0.5, f"分析主题：{meta['subtitle']}", 13, color=TEXT)
    _box(slide, 0.55, 2.65, 5.9, 0.5, _qa_line(status, passed, checks), 14,
         bold=True, color=OK if status == "PASS" else BAD)
    _box(slide, 0.55, 3.35, 5.9, 2.6,
         "分析链路：数据概况 → 指标计算（公式→代入→结果→解释）→ 图表 → 分层结论 → QA 三层验证",
         12, color=GRAY)

    # 右侧首张图表
    charts = sorted((config.OUTPUT / f"task{task_no:02d}" / "charts").glob("*.png"))
    if charts:
        pic = charts[0]
        try:
            from PIL import Image
            w, h = Image.open(pic).size
        except Exception:
            w, h = 1280, 720
        box_w, box_h = 5.9, 5.4
        ratio = min(box_w / w, box_h / h)
        pw, ph = w * ratio, h * ratio
        slide.shapes.add_picture(str(pic), Inches(7.0), Inches(1.5),
                                 width=Inches(pw), height=Inches(ph))
    else:
        _box(slide, 7.0, 3.0, 5.9, 0.5, "（无图表）", 13, color=GRAY)


def summary_pptx(out_path, log_path: Path | None = None) -> Path:
    """总览 PPTX：封面 + 能力 + 每任务一页 + 统计页。"""
    log_path = log_path or (config.OUTPUT / "build_log.json")
    log = {"tasks": []}
    if log_path.exists():
        import json
        log = json.loads(log_path.read_text(encoding="utf-8"))
    entries = {e["task"]: e for e in log["tasks"]}

    prs = _new_pres()

    # ---- 封面 ----
    s = _blank(prs)
    shape = s.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.35))
    shape.fill.solid(); shape.fill.fore_color.rgb = BRAND; shape.line.fill.background()
    _box(s, 1.0, 2.1, 11.3, 0.5, "AI DATA ANALYSIS", 16, bold=True, color=BRAND,
         align=PP_ALIGN.CENTER)
    _box(s, 1.0, 2.7, 11.3, 1.0, "AI 数据分析工作台 · 交付总览", 40, bold=True,
         color=BRAND_DARK, align=PP_ALIGN.CENTER)
    _box(s, 1.0, 3.8, 11.3, 0.5, "《统计与数据分析基础》课程配套 · 课堂演示工具", 16,
         color=GRAY, align=PP_ALIGN.CENTER)
    _box(s, 1.0, 6.6, 11.3, 0.4, "静态 HTML 工作台 · 双击即开 · 数据只读 · 结果可导出本 PPTX",
         12, color=GRAY, align=PP_ALIGN.CENTER)

    # ---- 能力介绍 ----
    s = _blank(prs)
    _bar(s, "工作台能力一览", "9 个实训任务全流程跑通 + 1 个外部工具占位")
    feats = [
        ("9 大实训全流程", "数据概况 → 指标计算 → 图表 → 分层结论 → QA 验证，报告即标准答案"),
        ("对齐课堂风格", "AI DATA ANALYSIS 标签 + DIG 框架 + 编号步骤卡 + 公式四件套，可直接投影"),
        ("推荐 Prompt 可复现", "每个任务附提示词，学生课后复制给 AI 独立复现同款分析"),
        ("数据安全", "原始数据只读（SHA-256 校验）、示例数据按课程描述构造、可替换真实数据重跑"),
        ("全离线", "单文件 HTML + 本地图表 PNG，不依赖任何外部服务"),
    ]
    for i, (t, x) in enumerate(feats):
        y = 1.6 + i * 1.05
        _box(s, 0.8, y, 1.9, 0.5, f"0{i + 1}", 20, bold=True, color=ACCENT)
        _box(s, 2.9, y, 2.6, 0.5, t, 16, bold=True, color=BRAND_DARK)
        _box(s, 5.8, y, 7.0, 0.9, x, 13, color=TEXT)

    # ---- 每任务一页 ----
    for task_no in sorted(config.TASKS):
        entry = entries.get(task_no, {"status": "?"})
        _task_slide(prs, task_no, entry)

    # ---- 统计页 ----
    s = _blank(prs)
    _bar(s, "构建统计", "output/build_log.json 全量记录")
    n_pass = sum(1 for e in log["tasks"] if e["status"] == "PASS")
    n_skip = sum(1 for e in log["tasks"] if e["status"] == "SKIP")
    _box(s, 0.8, 2.0, 11.5, 1.0, f"✅ 通过 {n_pass} / 9 个可构建任务", 30, bold=True,
         color=OK)
    _box(s, 0.8, 3.1, 11.5, 0.6, f"⛔ 占位 1 个（任务 02 · 需外部采集工具）", 16, color=GRAY)
    _box(s, 0.8, 4.0, 11.5, 0.6, "所有指标由确定性代码计算（pandas/numpy/scipy），LLM 不做心算",
         13, color=TEXT)
    _box(s, 0.8, 6.9, 11.5, 0.4, "本文件由 engine/ppt_out.py 自动生成，可在 PowerPoint / WPS 中自由编辑",
         11, color=GRAY)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def export_task(task_no: int, out_path) -> Path:
    """单任务交付页 PPTX（独立导出，便于单独发给学生）。"""
    log_path = config.OUTPUT / "build_log.json"
    entries = {}
    if log_path.exists():
        import json
        entries = {e["task"]: e for e in json.loads(log_path.read_text(encoding="utf-8"))["tasks"]}
    prs = _new_pres()
    _task_slide(prs, task_no, entries.get(task_no, {"status": "?"}))
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    p1 = summary_pptx(config.OUTPUT / "交付总览.pptx")
    print("已生成:", p1)
